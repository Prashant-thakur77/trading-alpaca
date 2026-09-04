#!/usr/bin/env python3
"""
Build the full reviewer presentation: one slide list -> PDF (via Chromium)
and .pptx (via python-pptx), from real screenshots in docs/press/deck-assets.

    python3 scripts/build_presentation.py

Every figure on a slide is one the repo can reproduce: the test count is read
at build time, the rest are the values verified in docs/press/video-script.md
and the journal. Nothing here is a projection.
"""
import os
import re
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "docs" / "press" / "deck-assets"
OUT_HTML = ROOT / "docs" / "press" / "presentation.html"
OUT_PDF = ROOT / "docs" / "press" / "trading-alpaca-presentation.pdf"
OUT_PPTX = ROOT / "docs" / "press" / "trading-alpaca-presentation.pptx"

SITE = "https://trading-alpaca-judge.vercel.app"
REPO = "https://github.com/Prashant-thakur77/trading-alpaca"
RAW = REPO + "/raw/main"

sys.path.insert(0, str(ROOT / "scripts"))
from build_deck import _playwright_node_path, tests  # noqa: E402


# ── image prep ────────────────────────────────────────────────
def crop(src: str, box, dst: str) -> str:
    """Crop with Pillow if available; otherwise use the original."""
    try:
        from PIL import Image
        im = Image.open(ASSETS / src)
        im.crop(box).save(ASSETS / dst)
        return dst
    except Exception:
        return src


HERO = crop("site-hero.png", (0, 560, 1660, 1080), "site-hero-crop.png")        # headline + CTA, no sidebar
INVERT = crop("site-stages.png", (300, 0, 1660, 300), "site-invert-crop.png")   # "Most agents ask..." headline
BUILD = crop("site-stages.png", (300, 500, 1660, 1080), "site-build-crop.png")  # BUILD stage + 3 numbers
REFUSE = crop("site-lower.png", (300, 150, 1660, 700), "site-refuse-crop.png")  # "IT REFUSED 72%"
CANDS = crop("judge-stages.png", (380, 0, 1540, 800), "judge-cands-crop.png")   # candidate table
SMILE = crop("smile.png", (330, 30, 1590, 1080), "smile-crop.png")
FILL = "term-fill.png"
# The committee block of the Tuesday fill cycle: frame at 36.5s of demo.mp4,
# cropped to the vol_analyst / bear_adversary / aggregate lines so the text is
# legible at slide scale. Extracted once with:
#   ffmpeg -ss 36.5 -i docs/press/demo.mp4 -frames:v 1 deck-assets/term-committee-full.png
COMMITTEE = crop("term-committee-full.png", (0, 60, 1920, 720), "term-committee.png")
CLOSE = crop("term-close.png", (0, 0, 1920, 460), "term-close-crop.png")
SUMMARY = crop("term-summary.png", (0, 0, 1300, 860), "term-summary-crop.png")
DENY = crop("judge-deny.png", (380, 540, 1540, 1000), "judge-deny-crop.png")
FAILCLOSED = crop("judge-failclosed.png", (380, 540, 1540, 1000), "judge-failclosed-crop.png")


# ── slides ────────────────────────────────────────────────────
# Each slide: tag, title, bullets (list[str]) or table (list[list]) or flow,
# optional image, optional note. Both renderers read this list.
SLIDES = [
    dict(tag="Alpaca AI Trading Agents · Options Alpha Agents",
         title="Trading Alpaca:\nthe options desk\nthat grades itself",
         bullets=[
             "Deterministic code builds every trade. A language model may only pick one by id, or refuse.",
             "Two independent vetoes · a fail-closed risk guard · a hash-chained journal · analysts that are Brier-scored and lose their vote when miscalibrated.",
             f"Live site: {SITE}",
             f"Judge desk (no credentials needed): {SITE}/judge",
             f"Code: {REPO}",
         ],
         image=HERO, layout="split"),

    dict(tag="The problem",
         title="Most agents ask a model what to buy.\nThis one never lets it answer.",
         bullets=[
             "Asking an LLM \"what should I trade?\" is fast to build and impossible to trust.",
             "It can hallucinate a strike, size a position wrongly, or be confidently wrong with no record of why.",
             "So this desk hands the model a numbered menu it did not write. It may return one id, or the word ABSTAIN.",
             "A hallucinated id is treated as an abstention. There is no code path by which a model can invent a strike, change a quantity, or move a price.",
         ],
         image=INVERT, layout="split"),

    dict(tag="How it works",
         title="Six stages. Four of them can say no.",
         flow=[
             ("BUILD", "code enumerates every legal defined-risk structure from the live chain, fully priced"),
             ("ARGUE", "a vol analyst and a bear adversary each return a probability; a trader picks an id or ABSTAINs"),
             ("VETO", "a pure-code thesis check, then a blind LLM review that never sees the committee's reasoning"),
             ("GUARD", "RiskGuard enforces risk.yaml: ALLOW / DENY / DOWNSIZE, fail-closed on any error"),
             ("EXECUTE", "one atomic multi-leg order via the Alpaca CLI; a pre-mortem compiles exit triggers"),
             ("GRADE", "closes resolve each analyst's predictions; Brier scores recompute their weights"),
         ],
         note="Every stage appends to a hash-chained journal. Every refusal is recorded with its reason."),

    dict(tag="Stage 1 · Build",
         title="Every candidate is fully specified\nbefore any model is called.",
         bullets=[
             "Python reads the live SPY option chain and constructs every legal defined-risk structure: bull put, bear call, bull call and bear put spreads, iron condors, long straddles and butterflies.",
             "Each has strikes, legs, quantity, limit price, max loss and breakeven fixed by code.",
             "Liquidity gates first: open interest ≥ 100, quote width ≤ 10% of mid, 7–45 days to expiry.",
             "A live cycle reads ~3,700 contracts, ~1,100 pass liquidity, ~600–1,800 structures are built, 12 are shown to the committee.",
             "There is no function in the codebase that produces a naked short option.",
         ],
         image=CANDS, layout="split"),

    dict(tag="Stage 2 · Argue",
         title="Probabilities, not verdicts.\nAn adversary argues against every trade.",
         bullets=[
             "vol_analyst judges implied vs realized volatility; bear_adversary's only job is to find what breaks the trade. Each returns a probability, and an abstaining analyst is excluded from both numerator and denominator.",
             "The trader picks one id or ABSTAINs. In this live cycle it moved from c1 to c2 because the adversary flagged c1's at-the-money short strike — the objection changed the trade, and is journalled either way.",
         ],
         image=COMMITTEE, layout="wide"),

    dict(tag="Stage 3 · Veto",
         title="Two reviewers with decorrelated failure modes.",
         bullets=[
             "Two calls to the same model on the same context agree with each other and prove nothing.",
             "Gate 1 is pure code: do the position's Greeks match the structure's own directional thesis? Unmeasurable Greeks fail closed.",
             "Gate 2 is an LLM deliberately starved of the committee's reasoning, shown only the trade's own facts.",
             "Live on Thursday: a candidate passed the thesis check and was vetoed by the blind review — \"requires conviction SPY won't rally 1.4% over 44 days.\" The desk abstained.",
         ],
         image=FAILCLOSED, layout="split"),

    dict(tag="Stage 4 · Guard",
         title="risk.yaml is the single source of truth.\nThe guard fails closed.",
         table=[["Limit", "Value"],
                ["Max loss per position", "$1,000 (1% of equity)"],
                ["Concurrent positions", "3"],
                ["New trades per underlying per day", "1"],
                ["Net delta / net vega", "|Δ| ≤ 30 · |ν| ≤ 200"],
                ["Days to expiry", "7 – 45"],
                ["Liquidity", "OI ≥ 100 · spread ≤ 10% of mid"],
                ["Daily loss halt", "2% of equity"],
                ["After 3 consecutive losers", "size halved"],
                ["Any error, missing data, or exception", "DENY"]],
         note=f"Verified by {tests()} tests, including mutation tests that prove a guard exception denies rather than trades. A kill-switch file halts everything within a minute."),

    dict(tag="Stage 5 · Execute and exit",
         title="It manages the book before it adds to it.",
         bullets=[
             "Orders are one atomic multi-leg limit through the official Alpaca CLI. Live fills came in at better than their limits ($2.35 vs $1.94, $1.48 vs $1.22).",
             "Before the trade is sent, a pre-mortem asks the model how it fails, then compiles those failure modes into deterministic exit triggers.",
             "Three exits apply to every position regardless of any model: 50% of credit captured, max loss, and a forced close at 3 DTE.",
             "Thursday: the breakeven trigger set at entry fired at 772.42 and the desk closed the spread itself at −$70.50, before the loss could grow toward its $352 max.",
         ],
         image=CLOSE, layout="wide"),

    dict(tag="Stage 6 · Grade itself",
         title="An analyst that is confidently wrong loses its vote.",
         table=[["Analyst", "Resolved", "Brier", "Weight", "Interpretation"],
                ["vol_analyst", "12", "0.162", "1.18", "beats the 0.25 always-0.5 baseline — upweighted"],
                ["bear_adversary", "12", "0.298", "0.90", "worse than baseline — down-weighted"]],
         note="Closes resolve the probabilities each analyst gave; Brier scores recompute weights every cycle and feed the next aggregate. Twelve resolved trades is a small sample and is stated as one — the loop is the point, not the score."),

    dict(tag="Refusal is the product",
         title="Over 43 replayed decisions,\nit refused 72% of the time.",
         bullets=[
             "The real committee replayed over 43 post-knowledge-cutoff windows: 31 refusals, 12 executable trades.",
             "Every refusal gave a reason. The reasons are what matter: after the structure menu was fixed, zero refusals cited a missing structure; seven of eight were the two reviewers failing to agree.",
             "Replayed at $0.00: every LLM call is content-addressed and cached with its prompt, raw response and parse.",
             "This is not timidity. A desk that cannot say no cannot be trusted when it says yes.",
         ],
         image=REFUSE, layout="split"),

    dict(tag="Analytics",
         title="The volatility smile — and its noise floor.",
         bullets=[
             "A cubic fitted per expiry on ln(K/S)/√T across the liquid strikes only; 14 of 17 liquid groups fit.",
             "A strike is coloured only if its distance from the curve exceeds its own bid-ask spread. Everything grey scores exactly zero.",
             "We measured it, and deliberately did not let it pick strikes: the median deviation is worth about $3.94 per contract against a $4.00 quoted spread.",
             "Measurement is not edge, and the desk knows the difference.",
         ],
         image=SMILE, layout="split"),

    dict(tag="Live · official measurement window",
         title="What actually happened.",
         table=[["", ""],
                ["Account", "PA3JR0GVVEN0 — new, dedicated, paper, created 2026-08-29"],
                ["Opening equity", "$100,000.00, flat until the window opened"],
                ["Closing equity (Thu 3 Sep EOD)", "$99,827.64   (−0.17%)"],
                ["Fills", "3, all inside the window, each at better than its limit"],
                ["Closes", "1, automated, on its own breakeven trigger (−$70.50)"],
                ["Guard refusals", "per-day cap, position cap, no-stacking"],
                ["Positions corrupted", "0"],
                ["Every trade", "a defined-risk bear call spread, 1 contract, inside its $1,000 cap"]],
         note="Judged on total equity, this is flat. Three bear call spreads were one directional view expressed three times, and SPY rallied. We state that rather than hide it."),

    dict(tag="Found by running, not by reading",
         title="Four bugs, all while the test suite was green.",
         table=[["Day", "Bug", "What it cost", "Fix"],
                ["Mon", "Skill leak", "the host's plugin told analysts to look for skills first → silent abstain", "--setting-sources \"\" isolates the model call"],
                ["Mon", "Unfillable limit", "mid-priced order rested 55 min, blocked every later cycle", "20% marketable concession at the one live pricing site"],
                ["Tue", "Trailing comma", "a valid, reasoned trade decision discarded over one character", "string-aware JSON repair, fail-closed otherwise"],
                ["Thu", "Leg collision", "new spread sold a contract already held long → broker rejected", "filter candidates against current holdings"]],
         note="Each was root-caused from the journal and prompt cache, fixed with a regression test, and shipped mid-session. That is the argument for the journal."),

    dict(tag="What this does not claim",
         title="Stated limits, not omitted ones.",
         bullets=[
             "One closed trade and two open ones prove nothing statistically, and we do not imply otherwise.",
             "Paper fills are simulated against live quotes. No counterparty took the other side.",
             "The desk used ~25% of its own risk budget: candidates are built at 1 contract and the guard only ever downsizes. Known, documented, unchanged during scoring.",
             "Walk-forward backtests (13 windows, 90-day fit / 30-day test) are labelled as backtests; AAPL loses (PF 0.83) and the report says so.",
             "The previous author's hardcoded performance claims were deleted, not adapted.",
         ]),

    dict(tag="Try it yourself",
         title="Everything here is verifiable without credentials.",
         links=[
             ("Live site", SITE),
             ("Judge desk — replay four real decisions, two of them refusals", SITE + "/judge"),
             ("Volatility smile", SITE + "/smile"),
             ("Source, tests, journal, and every prompt", REPO),
             ("Demo video (4:22)", RAW + "/docs/press/demo.mp4"),
         ],
         cmds=["git clone " + REPO,
               "cd trading-alpaca && make install",
               "make test            # " + tests() + " tests, offline",
               "make judge           # 4 judge scenarios, no keys",
               "make verify-journal  # hash chain intact"],
         note="Compliance: orders route through the official Alpaca CLI; market data through the official Alpaca MCP server with read-only toolsets (no trading toolset is exposed, so an LLM cannot place an order through it). alpaca-py is used for analysis only."),
]


# ── HTML / PDF ────────────────────────────────────────────────
CSS = """
*{margin:0;padding:0;box-sizing:border-box}
:root{--ground:#161818;--g2:#0b0c0c;--ink:#f9f4eb;--dim:rgba(249,244,235,.70);--muted:rgba(249,244,235,.45);
 --hair:rgba(249,244,235,.14);--coral:#e75d60;--lime:#d0ff7e;
 --sans:"Inter Tight",Helvetica,Arial,sans-serif;--mono:"JetBrains Mono",Menlo,monospace}
@page{size:1280px 720px;margin:0}
body{background:var(--ground);color:var(--ink);font-family:var(--sans);line-height:1.4}
.slide{width:1280px;height:720px;padding:52px 64px 44px;display:flex;flex-direction:column;position:relative;
 page-break-after:always;overflow:hidden;background:var(--ground)}
.slide:last-child{page-break-after:auto}
.tag{font-family:var(--mono);font-size:11px;letter-spacing:.2em;text-transform:uppercase;color:var(--muted)}
.tag::before{content:"[ "}.tag::after{content:" ]"}
h1{font-size:44px;font-weight:800;text-transform:uppercase;letter-spacing:-.03em;line-height:.95;margin:12px 0 18px;white-space:pre-line}
.title h1{font-size:64px;margin-top:16px}
.body{display:flex;gap:36px;flex:1;min-height:0}
.body.split>.txt{flex:0 0 46%}.body.split>.img{flex:1}
.body.wide{flex-direction:column;gap:14px}.body.wide>.img{flex:1;min-height:0}
ul{list-style:none;display:flex;flex-direction:column;gap:11px}
li{color:var(--dim);font-size:16.5px;padding-left:20px;position:relative;line-height:1.38}
li::before{content:"\\2192";position:absolute;left:0;color:var(--coral)}
li b{color:var(--ink);font-weight:600}
.img{background:var(--g2);border-radius:6px;overflow:hidden;border:1px solid var(--hair);display:flex;align-items:center;justify-content:center}
.img img{width:100%;height:100%;object-fit:contain;object-position:top}
table{width:100%;border-collapse:collapse;font-size:15.5px;margin-top:6px}
th{text-align:left;color:var(--muted);font-family:var(--mono);font-weight:400;font-size:11px;letter-spacing:.12em;
 text-transform:uppercase;padding:8px 10px 8px 0;border-bottom:1px solid var(--hair)}
td{padding:8.5px 10px 8.5px 0;border-bottom:1px solid var(--hair);color:var(--dim);vertical-align:top}
td:first-child{color:var(--ink);font-family:var(--mono);font-size:14px;white-space:nowrap}
.note{border-left:2px solid var(--coral);padding:10px 0 10px 16px;color:var(--dim);font-size:14.5px;margin-top:auto;line-height:1.4}
.flow{display:grid;grid-template-columns:repeat(6,1fr);gap:10px;flex:1;align-content:center}
.flow .box{background:var(--g2);border:1px solid var(--hair);border-radius:6px;padding:16px 14px;min-height:250px;position:relative}
.flow .box .n{font-family:var(--mono);font-size:11px;color:var(--muted)}
.flow .box .k{font-size:22px;font-weight:800;letter-spacing:-.02em;margin:6px 0 10px}
.flow .box .k.no{color:var(--coral)}
.flow .box p{font-size:13.5px;color:var(--dim);line-height:1.38}
.flow .box::after{content:"\\2192";position:absolute;right:-13px;top:44%;color:var(--coral);font-size:18px}
.flow .box:last-child::after{content:""}
.links{display:flex;flex-direction:column;gap:12px;margin-top:6px}
.links a{color:var(--ink);text-decoration:none;font-size:18px}
.links a span{display:block;font-family:var(--mono);font-size:14px;color:var(--lime)}
pre{font-family:var(--mono);font-size:13.5px;color:var(--dim);background:var(--g2);border:1px solid var(--hair);
 border-radius:6px;padding:14px 16px;margin-top:14px;line-height:1.6;white-space:pre-wrap;word-break:break-all}
.n{position:absolute;right:64px;bottom:26px;font-family:var(--mono);font-size:11px;color:var(--muted)}
.foot{position:absolute;left:64px;bottom:26px;font-family:var(--mono);font-size:11px;color:var(--muted)}
"""


def esc(s: str) -> str:
    return (s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def html_slide(i: int, s: dict) -> str:
    n = len(SLIDES)
    parts = [f'<section class="slide{" title" if i == 0 else ""}">',
             f'<p class="tag">{esc(s["tag"])}</p>', f"<h1>{esc(s['title'])}</h1>"]
    if "flow" in s:
        parts.append('<div class="flow">')
        for j, (k, p) in enumerate(s["flow"], 1):
            no = " no" if k in ("ARGUE", "VETO", "GUARD", "GRADE") else ""
            parts.append(f'<div class="box"><div class="n">0{j}</div><div class="k{no}">{esc(k)}</div><p>{esc(p)}</p></div>')
        parts.append("</div>")
    elif "table" in s:
        rows = s["table"]
        parts.append("<table><tr>" + "".join(f"<th>{esc(c)}</th>" for c in rows[0]) + "</tr>")
        for r in rows[1:]:
            parts.append("<tr>" + "".join(f"<td>{esc(c)}</td>" for c in r) + "</tr>")
        parts.append("</table>")
    elif "links" in s:
        parts.append('<div class="body split"><div class="txt"><div class="links">')
        for label, url in s["links"]:
            parts.append(f'<a href="{url}">{esc(label)}<span>{esc(url)}</span></a>')
        parts.append("</div></div><div class=\"txt\"><pre>" + esc("\n".join(s["cmds"])) + "</pre></div></div>")
    else:
        layout = s.get("layout", "split")
        parts.append(f'<div class="body {layout}">')
        parts.append("<div class=\"txt\"><ul>" + "".join(f"<li>{esc(b)}</li>" for b in s["bullets"]) + "</ul></div>")
        if s.get("image"):
            parts.append(f'<div class="img"><img src="deck-assets/{s["image"]}"></div>')
        parts.append("</div>")
    if s.get("note"):
        parts.append(f'<div class="note">{esc(s["note"])}</div>')
    parts.append(f'<div class="foot">trading-alpaca · paper account only · {SITE.replace("https://", "")}</div>')
    parts.append(f'<div class="n">{i + 1:02d} / {n:02d}</div></section>')
    return "\n".join(parts)


def build_html() -> str:
    head = ('<!doctype html><meta charset="utf-8"><title>Trading Alpaca — presentation</title>'
            '<link rel="preconnect" href="https://fonts.googleapis.com">'
            '<link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Inter+Tight:wght@400;600;800&family=JetBrains+Mono:wght@400;600&display=swap">'
            f"<style>{CSS}</style>")
    return head + "\n".join(html_slide(i, s) for i, s in enumerate(SLIDES))


def render_pdf() -> bool:
    mods = _playwright_node_path()
    if mods is None:
        print("  PDF render skipped: playwright not found"); return False
    js = f"""
const {{ chromium }} = require('playwright');
(async () => {{
  const b = await chromium.launch({{ executablePath:'/snap/bin/chromium', headless:true }});
  const p = await b.newPage();
  await p.goto('file://{OUT_HTML}', {{ waitUntil:'networkidle' }});
  await p.waitForTimeout(1500);
  await p.pdf({{ path:'{OUT_PDF}', width:'1280px', height:'720px', printBackground:true,
    margin:{{top:'0',bottom:'0',left:'0',right:'0'}} }});
  await b.close();
}})();
"""
    d = ROOT / ".pres-tmp"; d.mkdir(exist_ok=True); (d / "r.js").write_text(js)
    r = subprocess.run(["node", str(d / "r.js")], capture_output=True, text=True,
                       env={"NODE_PATH": str(mods), "PATH": "/usr/bin:/bin:/home/prashant/.volta/bin"})
    import shutil; shutil.rmtree(d, ignore_errors=True)
    if r.returncode == 0 and OUT_PDF.exists():
        print(f"  wrote {OUT_PDF.relative_to(ROOT)} ({OUT_PDF.stat().st_size:,} bytes)"); return True
    print(f"  PDF render failed: {(r.stderr or '')[:300]}"); return False


# ── PPTX ──────────────────────────────────────────────────────
def build_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    GROUND = RGBColor(0x16, 0x18, 0x18); INK = RGBColor(0xF9, 0xF4, 0xEB)
    DIM = RGBColor(0xB8, 0xB4, 0xAC); MUTED = RGBColor(0x8A, 0x87, 0x80)
    CORAL = RGBColor(0xE7, 0x5D, 0x60); LIME = RGBColor(0xD0, 0xFF, 0x7E); G2 = RGBColor(0x0B, 0x0C, 0x0C)

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height
    ML, MT = Inches(0.65), Inches(0.5)

    def bg(slide):
        f = slide.background.fill; f.solid(); f.fore_color.rgb = GROUND

    def text(slide, x, y, w, h, s, size, color=INK, bold=False, font="Calibri", align=None, mono=False):
        tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
        lines = s.split("\n")
        for k, ln in enumerate(lines):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = ln; r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "Consolas" if mono else font
            if align: p.alignment = align
        return tb

    def bullets(slide, x, y, w, h, items, size=14):
        tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
        for k, it in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            r1 = p.add_run(); r1.text = "→ "; r1.font.color.rgb = CORAL; r1.font.size = Pt(size); r1.font.bold = True
            r2 = p.add_run(); r2.text = it; r2.font.color.rgb = DIM; r2.font.size = Pt(size); r2.font.name = "Calibri"
            p.space_after = Pt(7)

    def picture(slide, name, x, y, w, h):
        path = ASSETS / name
        if not path.exists(): return
        from PIL import Image
        iw, ih = Image.open(path).size
        # fit inside box, keep aspect
        scale = min(w / iw, h / ih); pw, ph = int(iw * scale), int(ih * scale)
        slide.shapes.add_picture(str(path), x + int((w - pw) / 2), y, width=pw, height=ph)

    def table(slide, rows, x, y, w, size=11):
        nrow, ncol = len(rows), len(rows[0])
        shp = slide.shapes.add_table(nrow, ncol, x, y, w, Inches(0.36) * nrow)
        t = shp.table
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                c = t.cell(i, j); c.text = cell
                c.fill.solid(); c.fill.fore_color.rgb = G2 if i else GROUND
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(size if i else 9); r.font.color.rgb = INK if (i and j == 0) or i == 0 else DIM
                        r.font.name = "Consolas" if (j == 0 and i) or i == 0 else "Calibri"
                        if i == 0: r.font.color.rgb = MUTED
        return shp

    n = len(SLIDES)
    for i, s in enumerate(SLIDES):
        sl = prs.slides.add_slide(blank); bg(sl)
        text(sl, ML, MT, W - 2 * ML, Inches(0.3), f"[ {s['tag'].upper()} ]", 9, MUTED, mono=True)
        tsize = 40 if i == 0 else 28
        text(sl, ML, MT + Inches(0.3), W - 2 * ML, Inches(1.3), s["title"].upper(), tsize, INK, bold=True)
        top = MT + Inches(1.75)
        body_h = H - top - Inches(0.9)
        if "flow" in s:
            bw = (W - 2 * ML - Inches(0.15) * 5) / 6
            for j, (k, p) in enumerate(s["flow"]):
                x = ML + j * (bw + Inches(0.15))
                box = sl.shapes.add_shape(1, x, top, bw, Inches(3.4)); box.fill.solid(); box.fill.fore_color.rgb = G2
                box.line.color.rgb = RGBColor(0x33, 0x35, 0x35)
                text(sl, x + Inches(0.12), top + Inches(0.1), bw - Inches(0.24), Inches(0.3), f"0{j+1}", 9, MUTED, mono=True)
                text(sl, x + Inches(0.12), top + Inches(0.4), bw - Inches(0.24), Inches(0.5), k, 16, CORAL if k in ("ARGUE", "VETO", "GUARD", "GRADE") else INK, bold=True)
                text(sl, x + Inches(0.12), top + Inches(0.95), bw - Inches(0.24), Inches(2.3), p, 11, DIM)
        elif "table" in s:
            table(sl, s["table"], ML, top, W - 2 * ML, size=11 if len(s["table"]) < 8 else 10)
        elif "links" in s:
            y = top
            for label, url in s["links"]:
                text(sl, ML, y, Inches(7.2), Inches(0.35), label, 13, INK)
                text(sl, ML, y + Inches(0.3), Inches(7.2), Inches(0.3), url, 10, LIME, mono=True); y += Inches(0.72)
            text(sl, ML + Inches(7.4), top, Inches(5.3), Inches(2.4), "\n".join(s["cmds"]), 10, DIM, mono=True)
        else:
            if s.get("image") and s.get("layout") == "wide":
                bullets(sl, ML, top, W - 2 * ML, Inches(1.9), s["bullets"], size=12)
                picture(sl, s["image"], ML, top + Inches(2.0), W - 2 * ML, body_h - Inches(2.0))
            elif s.get("image"):
                tw = Inches(5.7)
                bullets(sl, ML, top, tw, body_h, s["bullets"], size=12.5)
                picture(sl, s["image"], ML + tw + Inches(0.3), top, W - 2 * ML - tw - Inches(0.3), body_h)
            else:
                bullets(sl, ML, top, W - 2 * ML, body_h, s["bullets"], size=14)
        if s.get("note"):
            nb = text(sl, ML + Inches(0.15), H - Inches(1.05), W - 2 * ML - Inches(0.15), Inches(0.6), s["note"], 10, DIM)
            bar = sl.shapes.add_shape(1, ML, H - Inches(1.05), Inches(0.03), Inches(0.6)); bar.fill.solid(); bar.fill.fore_color.rgb = CORAL; bar.line.fill.background()
        text(sl, ML, H - Inches(0.4), Inches(6), Inches(0.25), f"trading-alpaca · paper account only · {SITE.replace('https://','')}", 8, MUTED, mono=True)
        text(sl, W - ML - Inches(1.2), H - Inches(0.4), Inches(1.2), Inches(0.25), f"{i+1:02d} / {n:02d}", 8, MUTED, mono=True, align=PP_ALIGN.RIGHT)
    prs.save(OUT_PPTX)
    print(f"  wrote {OUT_PPTX.relative_to(ROOT)} ({OUT_PPTX.stat().st_size:,} bytes, {n} slides)")


def main() -> int:
    OUT_HTML.write_text(build_html())
    print(f"  wrote {OUT_HTML.relative_to(ROOT)} ({len(SLIDES)} slides)")
    ok = render_pdf()
    build_pptx()
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
